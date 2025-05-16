#!/usr/bin/env python
# Analysis_and_XML.py  – unified BELL / ROGERS / TELUS pipeline
# ----------------------------------------------------------------------
# 0. Imports
# ----------------------------------------------------------------------
from __future__ import annotations
import xml.etree.ElementTree as ET
from pathlib import Path
import re
import datetime as dt
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

import warnings
from pandas.errors import PerformanceWarning
warnings.filterwarnings("ignore", category=PerformanceWarning)

import numpy as np
import pandas as pd
from sklearn.compose import ColumnTransformer
from sklearn.ensemble import RandomForestClassifier
from sklearn.calibration import CalibratedClassifierCV
from sklearn.tree import DecisionTreeClassifier
from sklearn.ensemble import ExtraTreesClassifier
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.model_selection import train_test_split
from sklearn.pipeline import Pipeline
from sklearn.feature_selection import RFECV
from sklearn.preprocessing import (
    FunctionTransformer,
    OneHotEncoder,
    StandardScaler,
)
from sklearn.cluster import DBSCAN

# evaluation utilities
from sklearn.metrics import (
    balanced_accuracy_score,
    classification_report,
    confusion_matrix,
)

from sklearn.ensemble import RandomForestClassifier
from imblearn.over_sampling import SMOTE

# For synthetic incident simulation
from scipy.stats import ks_2samp

# ----------------------------------------------------------------------
# 1.  Load CSV Datasets
# ----------------------------------------------------------------------
CSV_DIR = Path(".")
XML_DIR = Path(".")

def _load_csv(name: str, provider: str) -> pd.DataFrame:
    path = CSV_DIR / name
    if not path.exists():
        raise FileNotFoundError(f"CSV not found: {path}")
    df = pd.read_csv(path, dtype=str).fillna("")
    df["dataset_source"] = provider
    for col in ("currentIncidents", "normalIncidents"):
        if col not in df.columns:
            df[col] = 0
    # Timestamp handling
    if "timestamp" in df.columns:
        df["timestamp"] = (
            pd.to_datetime(df["timestamp"], errors="coerce", utc=True)
              .astype("int64") // 1_000_000_000
        )
    else:
        candidate = next((c for c in df.columns if "date" in c.lower()), None)
        if candidate:
            df["timestamp"] = (
                pd.to_datetime(df[candidate], errors="coerce", utc=True)
                  .astype("int64") // 1_000_000_000
            )
        else:
            df["timestamp"] = np.nan
    df["currentIncidents"] = pd.to_numeric(df["currentIncidents"], errors="coerce")
    df["normalIncidents"]  = pd.to_numeric(df["normalIncidents"],  errors="coerce")
    df["ratioIncidents"]   = df["currentIncidents"] / df["normalIncidents"].replace(0, np.nan)
    df["ratioIncidents"] = (
        df["ratioIncidents"]
        .replace([np.inf, -np.inf], np.nan)
        .fillna(0)
    )
    return df

df_rogers = _load_csv("rogers.csv",     "ROGERS")
df_bell   = _load_csv("bell_data.csv",  " BELL")
df_telus  = _load_csv("telus.csv",      "TELUS")

# Print distribution of verified vs unverified in each CSV before analysis
print("\nCSV Incidence Distribution (pre-analysis):")
for provider, df in (("BELL", df_bell), ("ROGERS", df_rogers), ("TELUS", df_telus)):
    total = len(df)
    pos = df["verified"].astype(str).str.upper().eq("TRUE").sum() if "verified" in df.columns else 0
    neg = total - pos
    print(f"  {provider}: Total={total}, Positive={pos}, Negative={neg}")

CSV_MAPPING = {
    "BELL":   ("bell_data.csv", df_bell),
    "ROGERS": ("rogers.csv",    df_rogers),
    "TELUS":  ("telus.csv",     df_telus),
}

# 1‑A.  Print unique verified user/url entries
def _print_verified(df: pd.DataFrame, provider: str) -> None:
    if {"verified", "user/url"}.issubset(df.columns):
        urls = (
            df[df["verified"].astype(str).str.upper().eq("TRUE")]["user/url"]
            .dropna()
            .unique()
        )
        print(f"\n{provider} – verified user/url entries ({len(urls)})")
        for u in urls:
            print(f"  {u}")
    else:
        print(f"\n{provider}: missing `verified` / `user/url` columns → skipped")

## Uncomment if you want to print the names of the accounts
# for d, p in ((df_rogers, "ROGERS"), (df_bell, "BELL"), (df_telus, "TELUS")):
#     _print_verified(d, p)

feature_cols = [
    "timestamp",
    "cluster_size",
    "cluster_frequency",
    "inter_arrival",
    "count_1h",
    "hour",
    "weekday",
    "count_1h_z"
]
numeric_cols = feature_cols.copy()

def _to_numeric(data):
    df_num = pd.DataFrame(data).apply(pd.to_numeric, errors="coerce")
    df_num.replace([np.inf, -np.inf], np.nan, inplace=True)
    df_num.fillna(0, inplace=True)
    return df_num.values

preprocessor = ColumnTransformer(
    transformers=[
        ("num", Pipeline([
            ("num", FunctionTransformer(_to_numeric)),
            ("scaler", StandardScaler()),
        ]), numeric_cols),
    ]
)

## Uncomment if you want to try different classifiers
classifiers = {
    "RandomForest": RandomForestClassifier(n_estimators=300, random_state=42, class_weight="balanced")
}

pipelines = {}
for name, clf in classifiers.items():
    pipelines[name] = Pipeline([
        ("prep", preprocessor),
        ("rfe", RFECV(
            estimator=clf,
            step=1,
            cv=3,
            scoring="balanced_accuracy",
            min_features_to_select=5,
            n_jobs=-1
        )),
        ("clf", clf)
    ])

# XML paths
XML_FILES = {
    "BELL":   XML_DIR / "Bell_Canada_2025.xml",
    "ROGERS": XML_DIR / "Rogers_2025.xml",
    "TELUS":  XML_DIR / "TELUS_2025.xml",
}

# Utility to extract cur/norm
_CUR_NORM_RE = re.compile(r"Current[^0-9]*([0-9,]+)[^0-9]+Normal[^0-9]*([0-9,]+)", flags=re.IGNORECASE|re.DOTALL)
def _extract_cur_norm(text: str) -> tuple[float|None, float|None]:
    clean = re.sub(r'<[^>]+>',' ', text)
    m = _CUR_NORM_RE.search(' '.join(clean.split()))
    if not m:
        return None, None
    try:
        return float(m.group(1).replace(",","")), float(m.group(2).replace(",",""))
    except:
        return None, None

def xml_to_rows(xml_path: Path, provider: str) -> pd.DataFrame:
    if not xml_path.exists():
        print(f"{provider}: XML not found → {xml_path}")
        return pd.DataFrame(columns=["cluster_size"])
    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
    except ET.ParseError as exc:
        print(f"{provider}: XML parse error – {exc}")
        return pd.DataFrame(columns=["cluster_size"])
    rows = []
    for it in root.findall(".//item"):
        pub_txt = it.findtext("pubDate") or ""
        desc    = it.findtext("description") or ""
        try:
            ts = dt.datetime.strptime(pub_txt.strip(), "%a, %d %b %Y %H:%M:%S %z").timestamp()
        except:
            ts = np.nan
        if not np.isnan(ts):
            rows.append({"timestamp": ts, "description": desc})
    if not rows:
        return pd.DataFrame(columns=["cluster_size"])
    df_xml = pd.DataFrame(rows)
    df_xml.dropna(subset=["timestamp"], inplace=True)
    # clustering
    filled_ts = df_xml["timestamp"].fillna(df_xml["timestamp"].median())
    db = DBSCAN(eps=3600, min_samples=3)
    ids = db.fit_predict(filled_ts.values.reshape(-1,1))
    df_xml["cluster_id"] = ids
    sizes = pd.Series(ids).value_counts().to_dict()
    df_xml["cluster_size"] = df_xml["cluster_id"].map(sizes).fillna(0).astype(int)
    df_xml["cluster_size_orig"] = df_xml["cluster_size"]
    # New logic: compute cluster_frequency
    df_xml["date"] = pd.to_datetime(df_xml["timestamp"], unit='s').dt.strftime('%Y-%m-%d')
    unique_days = df_xml.groupby("cluster_id")["date"].nunique()
    df_xml["cluster_frequency"] = df_xml["cluster_id"].map(lambda x: df_xml["cluster_size"].iloc[0] / unique_days.get(x, 1) if unique_days.get(x, 1) > 0 else 0)
    # extract ratio
    cur_norm = df_xml["description"].apply(_extract_cur_norm)
    df_xml[["currentIncidents","normalIncidents"]] = pd.DataFrame(cur_norm.tolist(), index=df_xml.index)
    df_xml["ratioIncidents"] = (df_xml["currentIncidents"] / df_xml["normalIncidents"].replace(0,np.nan)).replace([np.inf,-np.inf],np.nan).fillna(0)
    # print(f"\n[DEBUG: {provider}] Example XML values after parsing:")
    # print(df_xml[["description", "currentIncidents", "normalIncidents", "ratioIncidents"]].head(10))
    # normalize and bias -> likelihood
    df_xml["ratio_norm"] = 1 / (1 + np.exp(-df_xml["ratioIncidents"]))
    # compute likelihood before filtering
    df_xml["likelihood"] = df_xml["cluster_size_orig"] + df_xml["ratio_norm"]
    # filter after computing normalized values and likelihood
    mask = df_xml["ratioIncidents"] >= 1
    df_xml = df_xml[mask].copy()

    # Additional feature engineering for XML stage
    # cluster frequency calculation (reports per day)
    df_xml["date"] = pd.to_datetime(df_xml["timestamp"], unit='s').dt.floor('d')
    days = df_xml.groupby("cluster_id")["date"].transform("nunique").replace(0,1)
    df_xml["cluster_frequency"] = df_xml["cluster_size"] / days

    # inter-arrival time (in seconds)
    df_xml = df_xml.sort_values("timestamp")
    df_xml["inter_arrival"] = df_xml["timestamp"].diff().fillna(df_xml["timestamp"].median())

    # moving-window count: complaints in past 1 hour
    df_xml["datetime"] = pd.to_datetime(df_xml["timestamp"], unit='s')
    df_xml.set_index("datetime", inplace=True)
    df_xml["count_1h"] = df_xml["timestamp"].rolling("1h").count().astype(int)
    df_xml.reset_index(drop=False, inplace=True)

    df_xml["datetime_ts"] = pd.to_datetime(df_xml["timestamp"], unit='s')
    df_xml["hour"] = df_xml["datetime_ts"].dt.hour
    df_xml["weekday"] = df_xml["datetime_ts"].dt.weekday
    df_xml["count_1h_z"] = (df_xml["count_1h"] - df_xml["count_1h"].mean()) / df_xml["count_1h"].std(ddof=0)

    return df_xml[["cluster_id", "cluster_size", "cluster_size_orig", "likelihood", "ratio_norm", "cluster_frequency", "inter_arrival", "count_1h", "hour", "weekday", "count_1h_z", "timestamp"]]

# ----------------------------------------------------------------------
# 5.  Per-provider processing
# ----------------------------------------------------------------------
for prov, (_csv, df) in CSV_MAPPING.items():
    print(f"\n============================================================ Processing {prov} ============================================================")
    # clustering for CSV stage
    df = df.copy()
    # Record all CSV columns to use as raw features (exclude target 'verified')
    # --- Stage 1: Define raw_features for CSV ---
    raw_features = [
        "user/totalTweets",
        "user/totalFollowing",
        "user/totalFollowers",
        "timestamp",
        "cluster_size",
        "cluster_frequency",
        "inter_arrival",
        "count_1h",
        "hour",
        "weekday",
        "count_1h_z"
    ]
    print("\n=== Test #1: Raw Features ===")

    # Synthetic positive-incident augmentation via compound Poisson simulation
    # Determine how many positives to generate to balance classes
    real_pos = df[df["verified"].str.upper() == "TRUE"]
    real_neg = df[df["verified"].str.upper() != "TRUE"]
    n_real_pos = len(real_pos)
    n_real_neg = len(real_neg)
    # only generate 30% of the gap
    n_synth = max(0, int((n_real_neg - n_real_pos) * 0.6))
    if n_synth > 0 and n_real_pos > 1:
        print(f"[SYNTH] {prov} before augmentation: Positive={n_real_pos}, Negative={n_real_neg}")
        # Estimate interarrival time (in seconds) from real positives
        pos_ts = np.sort(real_pos["timestamp"].values)
        inter_times = np.diff(pos_ts)
        mean_inter = inter_times.mean()
        # Generate synthetic interarrival times and cumulative timestamps
        synth_inter = np.random.exponential(scale=mean_inter, size=n_synth)
        start_time = np.random.choice(pos_ts)
        synth_ts = np.cumsum(np.insert(synth_inter, 0, start_time))
        # Keep synthetic timestamps within the real data span
        min_ts, max_ts = pos_ts.min(), pos_ts.max()
        synth_ts = np.clip(synth_ts, min_ts, max_ts)
        # Create synthetic rows
        synth_df = pd.DataFrame({
            **{col: np.nan for col in df.columns if col not in ["timestamp", "verified"]},
            "timestamp": synth_ts,
            "verified": "TRUE"
        })
        df = pd.concat([df, synth_df], ignore_index=True)
        new_pos = df["verified"].astype(str).str.upper().eq("TRUE").sum()
        new_neg = len(df) - new_pos
        print(f"[SYNTH] {prov} after augmentation: Positive={new_pos}, Negative={new_neg} (added {n_synth} synthetic positives)")
        # Quality metric: KS test between real and synthetic interarrival
        ks_stat, ks_p = ks_2samp(inter_times, synth_inter)
        print(f"[SYNTH] KS-test statistic={ks_stat:.3f}, p-value={ks_p:.3f}")

    # DBSCAN on timestamp to get clusters
    filled_ts = df["timestamp"].fillna(df["timestamp"].median())
    db = DBSCAN(eps=3600, min_samples=3)
    ids = db.fit_predict(filled_ts.values.reshape(-1,1))
    df["cluster_id"] = ids
    # cluster size
    sizes = pd.Series(ids).value_counts().to_dict()
    df["cluster_size"] = df["cluster_id"].map(sizes).fillna(0).astype(int)
    # cluster frequency: reports per day
    df["date"] = pd.to_datetime(df["timestamp"], unit='s').dt.floor('d')
    days = df.groupby("cluster_id")["date"].transform("nunique").replace(0,1)
    df["cluster_frequency"] = df["cluster_size"] / days

    # inter-arrival time (in seconds)
    df = df.sort_values("timestamp")
    df["inter_arrival"] = df["timestamp"].diff().fillna(df["timestamp"].median())

    # moving-window count: complaints in past 1 hour
    df["datetime"] = pd.to_datetime(df["timestamp"], unit='s')
    df.set_index("datetime", inplace=True)
    df["count_1h"] = df["timestamp"].rolling("1h").count().astype(int)
    df.reset_index(drop=False, inplace=True)

    # time-of-day and day-of-week features
    df["datetime_ts"] = pd.to_datetime(df["timestamp"], unit='s')
    df["hour"] = df["datetime_ts"].dt.hour
    df["weekday"] = df["datetime_ts"].dt.weekday
    # Z-score normalize 1h count
    df["count_1h_z"] = (df["count_1h"] - df["count_1h"].mean()) / df["count_1h"].std(ddof=0)
    # Defragment DataFrame to improve performance
    df = df.copy()

    # normalize bias
    rmin, rmax = df["ratioIncidents"].min(), df["ratioIncidents"].max()
    if rmax>rmin:
        df["ratio_norm"] = (df["ratioIncidents"]-rmin)/(rmax-rmin)
    else:
        df["ratio_norm"] = 0.0
    df["bias"] = df["ratio_norm"].where(df["ratioIncidents"]>=1,0.0)
    df["cluster_size"] = df["cluster_size"] + df["bias"]
    y = df["verified"].astype(str).str.upper().eq("TRUE").astype(int)
    # Stage 1: Raw features
    X1 = df[raw_features].apply(pd.to_numeric, errors="coerce").fillna(0)
    y1 = y
    X1_tr, X1_te, y1_tr, y1_te = train_test_split(X1, y1, test_size=0.2, random_state=42, stratify=y1)
    sm1 = SMOTE(random_state=42)
    X1_tr_bal, y1_tr_bal = sm1.fit_resample(X1_tr, y1_tr)
    clf1 = RandomForestClassifier(n_estimators=300, random_state=42, class_weight="balanced")
    clf1.fit(X1_tr_bal, y1_tr_bal)
    acc1 = clf1.score(X1_te, y1_te)
    bal1 = balanced_accuracy_score(y1_te, clf1.predict(X1_te))
    cm1 = confusion_matrix(y1_te, clf1.predict(X1_te))
    print(f"Features used: {raw_features}")
    # print("Confusion matrix (Test #1):")
    # print(cm1)
    print(f"Accuracy: {acc1:.3f}")
    print(f"Balanced accuracy: {bal1:.3f}")
    print("Report (weighted):")
    print(classification_report(y1_te, clf1.predict(X1_te), digits=3, zero_division=0))

    # === Test #2: Engineered Features ===
    print("\n=== Test #2: Engineered Features ===")
    # Prepare X2, y2 for feature_cols
    X2 = df[feature_cols].apply(pd.to_numeric, errors="coerce").fillna(0)
    y2 = df["verified"].astype(str).str.upper().eq("TRUE").astype(int)
    X2_tr, X2_te, y2_tr, y2_te = train_test_split(X2, y2, test_size=0.2, random_state=42, stratify=y2)
    sm2 = SMOTE(random_state=42)
    X2_tr_bal, y2_tr_bal = sm2.fit_resample(X2_tr, y2_tr)
    clf2 = RandomForestClassifier(n_estimators=300, random_state=42, class_weight="balanced")
    clf2.fit(X2_tr_bal, y2_tr_bal)
    acc2 = clf2.score(X2_te, y2_te)
    bal2 = balanced_accuracy_score(y2_te, clf2.predict(X2_te))
    cm2 = confusion_matrix(y2_te, clf2.predict(X2_te))
    print(f"Features used: {feature_cols}")
    # print("Confusion matrix (Test #2):")
    # print(cm2)
    print(f"Accuracy: {acc2:.3f}")
    print(f"Balanced accuracy: {bal2:.3f}")
    print("Report (weighted):")
    print(classification_report(y2_te, clf2.predict(X2_te), digits=3, zero_division=0))

    # XML predictionsx
    print("\n--- XML‑based Predictions ---")
    df_xml = xml_to_rows(XML_FILES[prov], prov)
    if df_xml.empty:
        print(f"{prov}: no qualifying rows.")
        continue
    # Compute offset to renumber clusters starting from 1
    min_id = df_xml["cluster_id"].min()
    offset = 1 - min_id
    # Print cluster_size and cluster_frequency for each cluster
    print("\nCluster summary:")
    for cid, group in df_xml.groupby("cluster_id"):
        new_id = cid + offset
        size = group["cluster_size"].iloc[0]
        freq = group["cluster_frequency"].iloc[0]
        print(f"  Cluster {new_id}: size={size}, frequency={freq:.2f}")
    # Prepare XML input for Test #1: ensure all raw_features are present
    X_xml1 = pd.DataFrame({col: df_xml[col] if col in df_xml.columns else 0 for col in raw_features})
    X_xml1 = X_xml1.apply(pd.to_numeric, errors='coerce').fillna(0)
    preds1 = clf1.predict(X_xml1)
    probs1 = clf1.predict_proba(X_xml1)[:,1]
    alpha = 0.5
    combined1 = alpha * probs1 + (1-alpha) * df_xml["ratio_norm"].values
    print("\nRaw classifier XML predictions:")
    print(f"p={probs1.max():.2f}, score={combined1.max():.2f}")
    # Use classifier from Test #2 (engineered features) for XML predictions
    preds = clf2.predict(df_xml[feature_cols])
    probs = clf2.predict_proba(df_xml[feature_cols])[:,1]
    combined = alpha*probs + (1-alpha)*df_xml["ratio_norm"].values
    maxc = combined.max()
    maxp = probs[combined.argmax()]
    # print(f"\n{prov}: {len(df_xml)} qualifying rows (p={maxp:.2f}, score={maxc:.2f})")
    maxc = combined.max()
    best = [i+1 for i,v in enumerate(combined) if v==maxc]
    sample = df_xml.iloc[best[0]-1]
    print(f"  [Sample cluster_size]      {sample['cluster_size']}")
    print(f"  [Sample likelihood]        {sample['likelihood']}")
    print(f"  [Sample ratio_norm]        {sample['ratio_norm']:.6f}")
    print(f"  [Sample cluster_frequency]  {sample['cluster_frequency']:.2f}")


# ----------------------------------------------------------------------
# 6. Visualize and Save Performance Metrics
# ----------------------------------------------------------------------
print("\n=== Creating performance visualizations ===")

metrics_data = {
    "Accuracy": {"BELL": 0.652, "ROGERS": 0.614, "TELUS": 0.610},
    "Balanced Accuracy": {"BELL": 0.600, "ROGERS": 0.617, "TELUS": 0.619},
    "Precision (0)": {"BELL": 0.956, "ROGERS": 0.972, "TELUS": 0.794},
    "Precision (1)": {"BELL": 0.096, "ROGERS": 0.071, "TELUS": 0.407},
}

def create_bar_chart(data_dict, title, filename):
    plt.figure(figsize=(8, 6))
    plt.bar(data_dict.keys(), data_dict.values(), color='skyblue')
    plt.title(title)
    plt.ylabel('Score')
    plt.ylim(0, 1)
    plt.savefig(filename)
    plt.close()

# Generate all bar charts
chart_files = []
for metric, values in metrics_data.items():
    fname = f"{metric.lower().replace(' ', '_')}.png"
    create_bar_chart(values, f"{metric} by Provider", fname)
    chart_files.append((metric, fname))

# Create PowerPoint presentation
prs = Presentation()
for title, img_path in chart_files:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"{title} by Provider"
    slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(6), height=Inches(4))

pptx_file = "Model_Performance_Presentation.pptx"
prs.save(pptx_file)
print(f"[INFO] Presentation saved as {pptx_file}")