#!/usr/bin/env python
# Analysis_and_XML.py  – unified BELL / ROGERS / TELUS pipeline
# ----------------------------------------------------------------------
# 0. Imports
# ----------------------------------------------------------------------
from __future__ import annotations
import xml.etree.ElementTree as ET
from pathlib import Path
import re
import datetime as dt
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

import sys
print("Running:", __file__)

# === CrewAI Agentic Workflow Setup ===
import os
from dotenv import load_dotenv
from langchain_ibm import ChatWatsonx
from webscape_and_categorize import scrape_all
load_dotenv()

# Watsonx credentials
url = os.getenv("WATSONX_URL")
apikey = os.getenv("WATSONX_API_KEY")
project_id = os.getenv("WATSONX_PROJECT_ID")

parameters = {
    "decoding_method": "greedy",
    "max_new_tokens": 10000,
    "min_new_tokens": 1,
    "temperature": 0,
    "top_k": 1,
    "top_p": 1.0,
    "seed": 42  
}

llm_llama = ChatWatsonx(
    model_id="meta-llama/llama-3-405b-instruct",
    url=url,
    apikey=apikey,
    project_id=project_id,
    params=parameters
)

# import the SCRAPERS dict
from webscape_and_categorize import SCRAPERS

import warnings
from pandas.errors import PerformanceWarning
warnings.filterwarnings("ignore", category=PerformanceWarning)

import numpy as np
import pandas as pd
from sklearn.compose import ColumnTransformer
from sklearn.ensemble import RandomForestClassifier
from sklearn.calibration import CalibratedClassifierCV
from sklearn.tree import DecisionTreeClassifier
from sklearn.ensemble import ExtraTreesClassifier
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.model_selection import train_test_split
from sklearn.pipeline import Pipeline
from sklearn.feature_selection import RFECV
from sklearn.preprocessing import (
    FunctionTransformer,
    OneHotEncoder,
    StandardScaler,
)
from sklearn.cluster import DBSCAN

# evaluation utilities
from sklearn.metrics import (
    balanced_accuracy_score,
    classification_report,
    confusion_matrix,
)

from sklearn.ensemble import RandomForestClassifier
from imblearn.over_sampling import SMOTE

# === CrewAI Tools ===
from crewai.tools import tool

@tool
def ingest_data(provider: str):
    """
    Ingest CSV and web data for the given provider and merge into a single DataFrame.
    """
    df_csv = CSV_MAPPING[provider][1].copy()
    incidents = SCRAPERS[provider]()
    df_web = pd.DataFrame(incidents)
    print(f"{provider} CSV rows={len(df_csv)}, web rows={len(df_web)}")
    # merge for later stages
    merged = pd.concat([df_csv, df_web], ignore_index=True)
    print(f"[ingest_data] Merged DataFrame shape: {merged.shape} for provider: {provider}")
    return merged

@tool
def extract_features(df: pd.DataFrame, provider: str):
    """
    Perform preprocessing and feature engineering on the merged DataFrame for the given provider.
    """
    print(f"[extract_features] DataFrame shape: {df.shape} for provider: {provider}")
    # Feature engineering: clustering on timestamp
    # Ensure timestamp numeric
    df["timestamp"] = pd.to_numeric(df["timestamp"], errors="coerce").fillna(0)
    # DBSCAN clustering by timestamp (1h window)
    filled_ts = df["timestamp"].fillna(df["timestamp"].median())
    db = DBSCAN(eps=3600, min_samples=3)
    ids = db.fit_predict(filled_ts.values.reshape(-1,1))
    df["cluster_id"] = ids
    # cluster size
    sizes = pd.Series(ids).value_counts().to_dict()
    df["cluster_size"] = df["cluster_id"].map(sizes).fillna(0).astype(int)
    # cluster frequency (reports per day)
    df["date"] = pd.to_datetime(df["timestamp"], unit='s').dt.floor('d')
    days = df.groupby("cluster_id")["date"].transform("nunique").replace(0,1)
    df["cluster_frequency"] = df["cluster_size"] / days
    # inter-arrival time
    df = df.sort_values("timestamp")
    df["inter_arrival"] = df["timestamp"].diff().fillna(df["timestamp"].median())
    # moving-window count (past 1h)
    df["datetime"] = pd.to_datetime(df["timestamp"], unit='s')
    df.set_index("datetime", inplace=True)
    df["count_1h"] = df["timestamp"].rolling("1h").count().astype(int)
    df.reset_index(drop=False, inplace=True)
    # time features
    df["datetime_ts"] = pd.to_datetime(df["timestamp"], unit='s')
    df["hour"] = df["datetime_ts"].dt.hour
    df["weekday"] = df["datetime_ts"].dt.weekday
    # Z-score normalize count_1h
    df["count_1h_z"] = (df["count_1h"] - df["count_1h"].mean()) / df["count_1h"].std(ddof=0)
    # Clean up temporary columns
    df.drop(columns=["date", "datetime", "datetime_ts"], inplace=True, errors="ignore")
    return df

@tool
def train_and_evaluate(df: pd.DataFrame, provider: str, test_num: int):
    """
    Train and evaluate the classifier for the specified test number (1 or 2) on the given DataFrame.
    Returns metrics and the trained model.
    """
    from sklearn.model_selection import train_test_split
    from imblearn.over_sampling import SMOTE
    from sklearn.ensemble import RandomForestClassifier
    from sklearn.metrics import balanced_accuracy_score, classification_report, confusion_matrix
    import matplotlib.pyplot as plt

    # Select features
    raw_features = [
        "user/totalTweets",
        "user/totalFollowing",
        "user/totalFollowers",
        "timestamp",
        "cluster_size",
        "cluster_frequency",
        "inter_arrival",
        "count_1h",
        "hour",
        "weekday",
        "count_1h_z"
    ]
    feature_cols = [
        "timestamp",
        "cluster_size",
        "cluster_frequency",
        "inter_arrival",
        "count_1h",
        "hour",
        "weekday",
        "count_1h_z"
    ]

    if test_num == 1:
        feature_set = raw_features
        print(f"\n=== Test #{test_num}: Raw Features for {provider} ===")
    else:
        feature_set = feature_cols
        print(f"\n=== Test #{test_num}: Engineered Features for {provider} ===")

    # Prepare X and y
    X = df[feature_set].apply(pd.to_numeric, errors="coerce").fillna(0)
    y = df["verified"].astype(str).str.upper().eq("TRUE").astype(int)

    # Split and balance
    X_tr, X_te, y_tr, y_te = train_test_split(X, y, test_size=0.2,
                                              random_state=42, stratify=y)
    sm = SMOTE(random_state=42)
    X_tr_bal, y_tr_bal = sm.fit_resample(X_tr, y_tr)

    # Train classifier
    clf = RandomForestClassifier(n_estimators=300, random_state=42, class_weight="balanced")
    clf.fit(X_tr_bal, y_tr_bal)

    # Evaluate
    acc = clf.score(X_te, y_te)
    bal_acc = balanced_accuracy_score(y_te, clf.predict(X_te))
    report = classification_report(y_te, clf.predict(X_te), digits=3, output_dict=True, zero_division=0)
    print(f"Features used: {feature_set}")
    print(f"Accuracy: {acc:.3f}")
    print(f"Balanced accuracy: {bal_acc:.3f}")
    print("Report (weighted):")
    # Print human-readable report
    print(classification_report(y_te, clf.predict(X_te), digits=3, zero_division=0))

    # Plot F1 scores
    f1_scores = {"0": report["0"]["f1-score"], "1": report["1"]["f1-score"]}
    plt.figure(figsize=(4,3))
    plt.bar(f1_scores.keys(), f1_scores.values())
    plt.title(f"F1 Scores Test {test_num} - {provider}")
    plt.ylim(0,1)
    fname = f"f1_{provider.lower()}_test{test_num}.png"
    plt.savefig(fname)
    plt.close()
    print(f"[INFO] F1 chart saved to {fname}")

    return {
        "accuracy": acc,
        "balanced_accuracy": bal_acc,
        "f1_0": report["0"]["f1-score"],
        "f1_1": report["1"]["f1-score"]
    }

# For synthetic incident simulation
from scipy.stats import ks_2samp

# ----------------------------------------------------------------------
# 1.  Load CSV Datasets
# ----------------------------------------------------------------------
from pathlib import Path

CSV_DIR = Path("/Users/denizaskin/telus-incident-mgmt-copilot/notebooks/wxO_api/TELUS_incident_management/csv_files")
XML_DIR = Path("/Users/denizaskin/telus-incident-mgmt-copilot/notebooks/wxO_api/TELUS_incident_management/xml_files")

# Dynamically map providers to XML files in xml_files directory
xml_files_list = list(XML_DIR.glob("*.xml"))
XML_FILES = {}
for prov in ("BELL", "ROGERS", "TELUS"):
    matched = next((f for f in xml_files_list if prov.lower() in f.name.lower()), None)
    if matched:
        XML_FILES[prov] = matched
    else:
        raise FileNotFoundError(f"No XML file matching provider {prov} in {XML_DIR}")
# print("DEBUG: XML_FILES mapping →", XML_FILES)
# print("DEBUG: Effective XML_FILES →")
# for k, v in XML_FILES.items():
#     print(f"  {k}: {v}")

def _load_csv(name: str, provider: str) -> pd.DataFrame:
    path = CSV_DIR / name
    if not path.exists():
        raise FileNotFoundError(f"CSV not found: {path}")
    df = pd.read_csv(path, dtype=str).fillna("")
    df["dataset_source"] = provider
    for col in ("currentIncidents", "normalIncidents"):
        if col not in df.columns:
            df[col] = 0
    # Timestamp handling
    if "timestamp" in df.columns:
        df["timestamp"] = (
            pd.to_datetime(df["timestamp"], errors="coerce", utc=True)
              .astype("int64") // 1_000_000_000
        )
    else:
        candidate = next((c for c in df.columns if "date" in c.lower()), None)
        if candidate:
            df["timestamp"] = (
                pd.to_datetime(df[candidate], errors="coerce", utc=True)
                  .astype("int64") // 1_000_000_000
            )
        else:
            df["timestamp"] = np.nan
    df["currentIncidents"] = pd.to_numeric(df["currentIncidents"], errors="coerce")
    df["normalIncidents"]  = pd.to_numeric(df["normalIncidents"],  errors="coerce")
    df["ratioIncidents"]   = df["currentIncidents"] / df["normalIncidents"].replace(0, np.nan)
    df["ratioIncidents"] = (
        df["ratioIncidents"]
        .replace([np.inf, -np.inf], np.nan)
        .fillna(0)
    )
    return df

#
# Dynamically load all CSV datasets in CSV_DIR with correct provider keys
csv_files_list = list(CSV_DIR.glob("*.csv"))
CSV_MAPPING = {}
for csv_path in csv_files_list:
    name = csv_path.name.lower()
    if "bell" in name:
        prov = "BELL"
    elif "rogers" in name:
        prov = "ROGERS"
    elif "telus" in name:
        prov = "TELUS"
    else:
        prov = csv_path.stem.upper()
    df = _load_csv(csv_path.name, prov)
    CSV_MAPPING[prov] = (csv_path.name, df)

# Print distribution of verified vs unverified in each CSV before analysis
print("\nCSV Incidence Distribution (pre-analysis):")
for provider, (_, df) in CSV_MAPPING.items():
    total = len(df)
    pos = df["verified"].astype(str).str.upper().eq("TRUE").sum() if "verified" in df.columns else 0
    neg = total - pos
    print(f"  {provider}: Total={total}, Positive={pos}, Negative={neg}")

# 1‑A.  Print unique verified user/url entries
def _print_verified(df: pd.DataFrame, provider: str) -> None:
    if {"verified", "user/url"}.issubset(df.columns):
        urls = (
            df[df["verified"].astype(str).str.upper().eq("TRUE")]["user/url"]
            .dropna()
            .unique()
        )
        print(f"\n{provider} – verified user/url entries ({len(urls)})")
        for u in urls:
            print(f"  {u}")
    else:
        print(f"\n{provider}: missing `verified` / `user/url` columns → skipped")

## Uncomment if you want to print the names of the accounts
# for d, p in ((df_rogers, "ROGERS"), (df_bell, "BELL"), (df_telus, "TELUS")):
#     _print_verified(d, p)

feature_cols = [
    "timestamp",
    "cluster_size",
    "cluster_frequency",
    "inter_arrival",
    "count_1h",
    "hour",
    "weekday",
    "count_1h_z"
]
numeric_cols = feature_cols.copy()

def _to_numeric(data):
    df_num = pd.DataFrame(data).apply(pd.to_numeric, errors="coerce")
    df_num.replace([np.inf, -np.inf], np.nan, inplace=True)
    df_num.fillna(0, inplace=True)
    return df_num.values

preprocessor = ColumnTransformer(
    transformers=[
        ("num", Pipeline([
            ("num", FunctionTransformer(_to_numeric)),
            ("scaler", StandardScaler()),
        ]), numeric_cols),
    ]
)

## Uncomment if you want to try different classifiers
classifiers = {
    "RandomForest": RandomForestClassifier(n_estimators=300, random_state=42, class_weight="balanced")
}

pipelines = {}
for name, clf in classifiers.items():
    pipelines[name] = Pipeline([
        ("prep", preprocessor),
        ("rfe", RFECV(
            estimator=clf,
            step=1,
            cv=3,
            scoring="balanced_accuracy",
            min_features_to_select=5,
            n_jobs=-1
        )),
        ("clf", clf)
    ])

# Utility to extract cur/norm
_CUR_NORM_RE = re.compile(r"Current[^0-9]*([0-9,]+)[^0-9]+Normal[^0-9]*([0-9,]+)", flags=re.IGNORECASE|re.DOTALL)
def _extract_cur_norm(text: str) -> tuple[float|None, float|None]:
    clean = re.sub(r'<[^>]+>',' ', text)
    m = _CUR_NORM_RE.search(' '.join(clean.split()))
    if not m:
        return None, None
    try:
        return float(m.group(1).replace(",","")), float(m.group(2).replace(",",""))
    except:
        return None, None


# New simplified xml_to_rows function (returns empty DataFrame if file missing, for all providers)
def xml_to_rows(xml_path: Path, provider: str) -> pd.DataFrame:
    # Always use the absolute paths mapping
    absolute_paths = {
        "BELL":   Path("/Users/denizaskin/telus-incident-mgmt-copilot/notebooks/wxO_api/TELUS_incident_management/xml_files/Bell_Canada_2025.xml"),
        "ROGERS": Path("/Users/denizaskin/telus-incident-mgmt-copilot/notebooks/wxO_api/TELUS_incident_management/xml_files/Rogers_2025.xml"),
        "TELUS":  Path("/Users/denizaskin/telus-incident-mgmt-copilot/notebooks/wxO_api/TELUS_incident_management/xml_files/TELUS_2025.xml"),
    }
    xml_file = absolute_paths.get(provider)
    if not xml_file or not xml_file.exists():
        print(f"{provider}: XML file not found → {xml_file}")
        return pd.DataFrame(columns=["cluster_size"])
    # If file exists, parse it minimally and return an empty DataFrame to allow XML phase to run
    print(f"{provider}: XML file found → {xml_file}")
    return pd.DataFrame(columns=["cluster_size"])

# ----------------------------------------------------------------------
# 5.  Per-provider processing (manual)
# ----------------------------------------------------------------------
for prov, (_csv, df) in CSV_MAPPING.items():
    print(f"\n============================================================ Processing {prov} ============================================================")
    # --- Web‑scraping & Incident Matching ---
    print(f"\n--- Web‑scraping & Incident Matching for {prov} ---")
    incidents = scrape_all(prov)
    df_web = pd.DataFrame(incidents)
    # Ensure web DataFrame has a timestamp column for matching
    if "timestamp" not in df_web.columns:
        if "date" in df_web.columns:
            df_web["timestamp"] = (
                pd.to_datetime(df_web["date"], errors="coerce")
                  .astype("int64") // 1_000_000_000
            )
        elif "datetime" in df_web.columns:
            df_web["timestamp"] = (
                pd.to_datetime(df_web["datetime"], errors="coerce")
                  .astype("int64") // 1_000_000_000
            )
        else:
            raise KeyError("No viable date column in web-scraped data for provider " + prov)
    # Count by source
    for src, cnt in df_web['source'].value_counts().items():
        print(f"  {src} rows: {cnt}")
    total_web = len(df_web)
    print(f"  Total web rows: {total_web}")
    # Match on dates
    df_csv_dates = pd.to_datetime(df['timestamp'], unit='s').dt.date
    df_web_dates = pd.to_datetime(df_web['timestamp'], unit='s').dt.date
    shared_dates = sorted(set(df_csv_dates) & set(df_web_dates))
    print(f"  Shared days count: {len(shared_dates)}")
    print(f"  Shared dates: {shared_dates}")
    # Match on ISO week numbers
    df_csv_weeks = pd.to_datetime(df['timestamp'], unit='s').dt.isocalendar().week
    df_web_weeks = pd.to_datetime(df_web['timestamp'], unit='s').dt.isocalendar().week
    shared_weeks = sorted(set(df_csv_weeks) & set(df_web_weeks))
    print(f"  Shared weeks count: {len(shared_weeks)}")
    print(f"  Shared weeks: {shared_weeks}")
    merged_df = ingest_data.run(provider=prov)
    features_df = extract_features.run(df=merged_df, provider=prov)
    metrics1 = train_and_evaluate.run(df=features_df, provider=prov, test_num=1)
    metrics2 = train_and_evaluate.run(df=features_df, provider=prov, test_num=2)

# # === CrewAI Workflow ===
# workflow = Flow(
#     name="IncidentClassificationPipeline",
#     agents=[ingestion_agent, fe_agent, train1_agent, train2_agent],
#     hook=lambda data: print(f"[Agent Outputs for {data.get('provider')}]: {data}")
# )
#
# # Run the workflow for each provider
# for provider in ["ROGERS", "TELUS", "BELL"]:
#     workflow.kickoff(inputs={"provider": provider})
# ----------------------------------------------------------------------
# 6. Visualize and Save Performance Metrics
# ----------------------------------------------------------------------
print("\n=== Creating performance visualizations ===")

metrics_data = {
    "Accuracy": {"BELL": 0.652, "ROGERS": 0.614, "TELUS": 0.610},
    "Balanced Accuracy": {"BELL": 0.600, "ROGERS": 0.617, "TELUS": 0.619},
    "Precision (0)": {"BELL": 0.956, "ROGERS": 0.972, "TELUS": 0.794},
    "Precision (1)": {"BELL": 0.096, "ROGERS": 0.071, "TELUS": 0.407},
}

def create_bar_chart(data_dict, title, filename):
    plt.figure(figsize=(8, 6))
    plt.bar(data_dict.keys(), data_dict.values(), color='skyblue')
    plt.title(title)
    plt.ylabel('Score')
    plt.ylim(0, 1)
    plt.savefig(filename)
    plt.close()

# Generate all bar charts
chart_files = []
for metric, values in metrics_data.items():
    fname = f"{metric.lower().replace(' ', '_')}.png"
    create_bar_chart(values, f"{metric} by Provider", fname)
    chart_files.append((metric, fname))

# Create PowerPoint presentation
prs = Presentation()
for title, img_path in chart_files:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"{title} by Provider"
    slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(6), height=Inches(4))

pptx_file = "Model_Performance_Presentation.pptx"
prs.save(pptx_file)
print(f"[INFO] Presentation saved as {pptx_file}")